from io import BytesIO

from fastapi.testclient import TestClient
from pptx import Presentation


def create_project(client: TestClient, document: dict) -> dict:
    return client.post(
        "/api/projects",
        json={"title": "Editable deck", "document_json": document},
    ).json()["data"]


def test_editable_pptx_export_and_import(client: TestClient, valid_document: dict) -> None:
    page = valid_document["pages"][0]
    valid_document["nodes"]["title"] = {
        "id": "title",
        "type": "text",
        "parentId": page["id"],
        "role": "title",
        "transform": {
            "x": 80,
            "y": 70,
            "width": 700,
            "height": 100,
            "rotation": 0,
            "scaleX": 1,
            "scaleY": 1,
        },
        "style": {
            "opacity": 1,
            "visible": True,
            "locked": False,
            "fill": {"type": "solid", "color": "#18181b"},
        },
        "text": {
            "content": "Editable title",
            "fontFamily": "Arial",
            "fontSize": 44,
            "fontWeight": 400,
            "lineHeight": 1.2,
            "letterSpacing": 0,
            "textAlign": "left",
            "runs": [{"start": 0, "end": 8, "fontWeight": 700}],
        },
    }
    page["children"] = ["title"]
    project = create_project(client, valid_document)

    exported = client.get(f"/api/pptx/export/{project['id']}")
    assert exported.status_code == 200
    deck = Presentation(BytesIO(exported.content))
    assert deck.slides[0].shapes[0].text == "Editable title"

    imported = client.post(
        "/api/pptx/import",
        files={
            "file": (
                "roundtrip.pptx",
                exported.content,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert imported.status_code == 201
    document = imported.json()["data"]["document_json"]
    assert document["pages"][0]["children"]
    assert any(node["type"] == "text" for node in document["nodes"].values())


def test_pptx_import_rejects_wrong_mime(client: TestClient) -> None:
    response = client.post(
        "/api/pptx/import",
        files={"file": ("bad.txt", b"bad", "text/plain")},
    )
    assert response.status_code == 400
